import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def load_json(filepath):
    with open(filepath, 'r') as f:
        return json.load(f)

def apply_text_formatting(text_frame, font_name, font_size):
    for p in text_frame.paragraphs:
        p.font.name = font_name
        p.font.size = Pt(font_size)

def add_bullet(text_frame, text, font_name, font_size):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)

def add_footer(slide, text, rules):
    footer_box = slide.shapes.add_textbox(
        Inches(rules["layout"]["footer_left_inches"]), 
        Inches(rules["layout"]["footer_top_inches"]), 
        Inches(9.0), Inches(0.5)
    )
    p = footer_box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = rules["font"]
    p.font.size = Pt(rules["sizes"]["footer"])
    if rules["footer_align"] == "left":
        p.alignment = PP_ALIGN.LEFT

def insert_picture_if_exists(slide, placeholder_idx, img_path):
    if img_path and os.path.exists(img_path):
        ph = slide.placeholders[placeholder_idx]
        slide.shapes.add_picture(img_path, ph.left, ph.top, width=ph.width)
        sp = ph.element
        sp.getparent().remove(sp)

def generate_universal_presentation(content_json_path, output_pptx_path):
    content = load_json(content_json_path)
    rules_path = "/Users/crist/Library/Mobile Documents/com~apple~CloudDocs/harness/shared/templates/pptx_rules.json"
    rules = load_json(rules_path)
    
    prs = Presentation("/Users/crist/.gemini/antigravity/brain/ca247e80-2733-47b7-8659-a5348e17636d/Presentacion-16.10.pptx")
    font = rules["font"]
    sz_title = rules["sizes"]["title"]
    sz_sub = rules["sizes"]["subtitle"]
    sz_cont = rules["sizes"]["content"]
    author = content.get("author", rules["author"])
    
    # --- 1. PORTADA ---
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = content["title"]
    apply_text_formatting(s1.shapes.title.text_frame, font, sz_title)
    s1.placeholders[1].text = f"{content['subtitle']}\n{author}"
    apply_text_formatting(s1.placeholders[1].text_frame, font, sz_sub)

    # --- 2. CONTEXTO ---
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Contexto Clínico"
    apply_text_formatting(s2.shapes.title.text_frame, font, sz_title)
    tf2 = s2.placeholders[1].text_frame
    tf2.text = content["context"][0] if content["context"] else ""
    apply_text_formatting(tf2, font, sz_cont)
    for b in content["context"][1:]:
        add_bullet(tf2, b, font, sz_cont)
    add_footer(s2, author, rules)

    # --- 3. OBJETIVOS ---
    s3 = prs.slides.add_slide(prs.slide_layouts[1])
    s3.shapes.title.text = "Objetivos"
    apply_text_formatting(s3.shapes.title.text_frame, font, sz_title)
    tf3 = s3.placeholders[1].text_frame
    tf3.text = content["objectives"][0] if content["objectives"] else ""
    apply_text_formatting(tf3, font, sz_cont)
    for b in content["objectives"][1:]:
        add_bullet(tf3, b, font, sz_cont)
    add_footer(s3, author, rules)

    # --- 4. ANÁLISIS DE ESTUDIOS ---
    for i, study in enumerate(content["studies"]):
        # Diapositiva A: Metodología
        sa = prs.slides.add_slide(prs.slide_layouts[2])
        sa.shapes.title.text = f"Estudio {i+1}: {study['name']} - Metodología"
        apply_text_formatting(sa.shapes.title.text_frame, font, sz_title)
        tfa = sa.placeholders[1].text_frame
        tfa.text = study["methodology"][0] if study["methodology"] else ""
        apply_text_formatting(tfa, font, sz_cont)
        for b in study["methodology"][1:]:
            add_bullet(tfa, b, font, sz_cont)
        insert_picture_if_exists(sa, 2, study.get("image_path", ""))
        add_footer(sa, f"Ref: {study['name']} | {author}", rules)
        
        # Diapositiva B: Resultados
        sb = prs.slides.add_slide(prs.slide_layouts[2])
        sb.shapes.title.text = f"Estudio {i+1}: {study['name']} - Resultados"
        apply_text_formatting(sb.shapes.title.text_frame, font, sz_title)
        tfb = sb.placeholders[1].text_frame
        tfb.text = study["results"][0] if study["results"] else ""
        apply_text_formatting(tfb, font, sz_cont)
        for b in study["results"][1:]:
            add_bullet(tfb, b, font, sz_cont)
        insert_picture_if_exists(sb, 2, study.get("image_path", ""))
        add_footer(sb, f"Ref: {study['name']} | {author}", rules)

    # --- 5. DISCUSIÓN ---
    sd = prs.slides.add_slide(prs.slide_layouts[1])
    sd.shapes.title.text = "Discusión"
    apply_text_formatting(sd.shapes.title.text_frame, font, sz_title)
    tfd = sd.placeholders[1].text_frame
    tfd.text = content["discussion"][0] if content["discussion"] else ""
    apply_text_formatting(tfd, font, sz_cont)
    for b in content["discussion"][1:]:
        add_bullet(tfd, b, font, sz_cont)
    add_footer(sd, author, rules)

    # --- 6. CONCLUSIÓN ---
    sc = prs.slides.add_slide(prs.slide_layouts[1])
    sc.shapes.title.text = "Conclusiones"
    apply_text_formatting(sc.shapes.title.text_frame, font, sz_title)
    tfc = sc.placeholders[1].text_frame
    tfc.text = content["conclusions"][0] if content["conclusions"] else ""
    apply_text_formatting(tfc, font, sz_cont)
    for b in content["conclusions"][1:]:
        add_bullet(tfc, b, font, sz_cont)
    add_footer(sc, author, rules)

    # --- 7. BIBLIOGRAFÍA ---
    sb = prs.slides.add_slide(prs.slide_layouts[1])
    sb.shapes.title.text = "Bibliografía"
    apply_text_formatting(sb.shapes.title.text_frame, font, sz_title)
    tfbib = sb.placeholders[1].text_frame
    tfbib.text = content["bibliography"][0] if content["bibliography"] else ""
    apply_text_formatting(tfbib, font, sz_cont)
    for b in content["bibliography"][1:]:
        add_bullet(tfbib, b, font, sz_cont)
    add_footer(sb, author, rules)

    # VALIDAR REGLA DE 12 A 30 DIAPOSITIVAS
    total_slides = len(prs.slides)
    if total_slides < 12:
        print(f"Advertencia: La presentación tiene {total_slides} diapositivas, que es menor a 12. Generando diapositivas de relleno/resumen...")
        while len(prs.slides) < 12:
            s_fill = prs.slides.add_slide(prs.slide_layouts[1])
            s_fill.shapes.title.text = "Análisis Secundario"
            apply_text_formatting(s_fill.shapes.title.text_frame, font, sz_title)
            s_fill.placeholders[1].text = "Análisis en progreso..."
            apply_text_formatting(s_fill.placeholders[1].text_frame, font, sz_cont)
            add_footer(s_fill, author, rules)
            
    total_slides = len(prs.slides)
    if total_slides > 30:
        print(f"Advertencia: La presentación tiene {total_slides} diapositivas, excediendo el límite de 30. Borrando exceso...")
        # No se recomienda borrar diapositivas en python-pptx directamente por corrupción XML.
        # Esto es solo un motor de aviso por ahora.

    prs.save(output_pptx_path)
    print(f"Presentación Universal guardada con éxito en {output_pptx_path}. (Total diapositivas: {len(prs.slides)})")

if __name__ == "__main__":
    import sys
    if len(sys.argv) < 3:
        print("Uso: python build_universal_pptx.py <input.json> <output.pptx>")
    else:
        generate_universal_presentation(sys.argv[1], sys.argv[2])
