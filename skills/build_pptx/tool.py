"""Skill build_pptx — materializa una presentación desde un esquema de slides."""
from __future__ import annotations
from pathlib import Path


def build_pptx(args: dict) -> str:
    title = args.get("title", "Presentación")
    slides = args.get("slides", [])
    out = args.get("out", "salida.pptx")
    try:
        from pptx import Presentation  # import perezoso
        from pptx.util import Pt
    except ImportError:
        return "ERROR: instala python-pptx (pip install python-pptx)."
    prs = Presentation()
    # Portada
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = title
    # Contenido
    for s in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = str(s.get("title", ""))
        body = slide.placeholders[1].text_frame
        body.clear()
        for i, b in enumerate(s.get("bullets", [])):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = str(b)
            p.font.size = Pt(18)
    Path(out).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return f"OK: presentación creada en {out} ({len(slides)} slides)."


def register_skill(reg) -> None:
    reg.register(
        "build_pptx",
        "Crea un .pptx desde {title, slides:[{title, bullets}], out}.",
        {"type": "object",
         "properties": {"title": {"type": "string"},
                        "slides": {"type": "array"},
                        "out": {"type": "string"}},
         "required": ["title", "slides", "out"]},
        build_pptx,
    )
